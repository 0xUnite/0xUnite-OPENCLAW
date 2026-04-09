from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.chart import XL_DATA_LABEL_POSITION

OUTPUT = '/Users/sudi/.openclaw/workspace/韩国TikTokShop从0到1指南_v2.0.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(247, 248, 252)
NAVY = RGBColor(25, 36, 74)
BLUE = RGBColor(45, 102, 245)
TEAL = RGBColor(16, 163, 127)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
GRAY = RGBColor(90, 98, 114)
LIGHT = RGBColor(255, 255, 255)
BORDER = RGBColor(217, 222, 233)
PURPLE = RGBColor(109, 40, 217)

FONT = 'PingFang SC'
FONT_EN = 'Aptos'


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(8.8), Inches(0.42))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = LIGHT
    if subtitle:
        st = slide.shapes.add_textbox(Inches(9.25), Inches(0.16), Inches(3.5), Inches(0.34))
        p2 = st.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = FONT_EN
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(210, 219, 255)


def add_footer(slide, text='韩国TikTok Shop从0到1指南 v2.0 · 仅供商业决策参考，最终以 TikTok / Amazon / 税务监管官方规则为准'):
    tb = slide.shapes.add_textbox(Inches(0.35), Inches(7.08), Inches(12.5), Inches(0.24))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(8)
    r.font.color.rgb = GRAY


def add_bullet_box(slide, left, top, width, height, title, bullets, fill_color=LIGHT, title_color=NAVY, accent=BLUE, font_size=16):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.12), height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.08), width - Inches(0.3), Inches(0.32))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = title_color

    content = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.45), width - Inches(0.35), height - Inches(0.5))
    tf = content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(40, 46, 56)
        p.space_after = Pt(4)
        p.bullet = True
    return shape


def add_text(slide, left, top, width, height, text, size=16, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def add_cover(slide):
    set_bg(slide, RGBColor(242, 246, 255))
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = RGBColor(240, 245, 255); band.line.fill.background()
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.75), Inches(0), Inches(3.58), Inches(7.5))
    side.fill.solid(); side.fill.fore_color.rgb = NAVY; side.line.fill.background()
    bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.0), Inches(8.2), Inches(3.15))
    bubble.fill.solid(); bubble.fill.fore_color.rgb = LIGHT; bubble.line.color.rgb = RGBColor(221, 229, 244)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(0.75), Inches(1.6), Inches(0.42))
    pill.fill.solid(); pill.fill.fore_color.rgb = ORANGE; pill.line.fill.background()
    add_text(slide, Inches(1.1), Inches(0.82), Inches(1.3), Inches(0.24), 'v2.0', 18, LIGHT, True)
    add_text(slide, Inches(1.0), Inches(1.35), Inches(7.4), Inches(1.6), '韩国主体切入TikTok Shop：\n韩国本地机会 + 东南亚跨境路径', 24, NAVY, True)
    add_text(slide, Inches(1.02), Inches(3.08), Inches(7.2), Inches(0.8), '把“韩国本土闭环”与“韩国主体做SEA跨境”彻底拆开讲。\n重点回答：谁能做、先做哪条路、怎么避坑、如何低成本练手。', 15, GRAY)
    add_text(slide, Inches(1.0), Inches(4.3), Inches(4.8), Inches(1.2), '适用对象\n• 韩国公司/韩国常住华人\n• 想做 TikTok Shop 或跨境内容电商的新手\n• 需要备选平台练手的人', 16, NAVY, True)
    add_text(slide, Inches(10.15), Inches(1.0), Inches(2.5), Inches(2.0), '修订重点', 22, LIGHT, True)
    add_text(slide, Inches(10.15), Inches(1.55), Inches(2.35), Inches(4.8), '1. 路径分流\n2. 本地站谨慎披露\n3. SEA跨境作为主路径\n4. VAT与合规加强免责声明\n5. 新手品类/决策树/亚马逊备选新增', 15, RGBColor(229, 235, 255))
    add_text(slide, Inches(10.15), Inches(6.55), Inches(2.3), Inches(0.4), '2026 / 实操修订版', 11, RGBColor(197, 206, 237), True)


def add_flow_cards(slide, title, cards, top=Inches(1.35), h=Inches(2.15)):
    gap = Inches(0.25)
    total_w = Inches(12.0)
    w = (total_w - gap * (len(cards)-1)) / len(cards)
    left = Inches(0.6)
    add_title(slide, title)
    for idx, (head, lines, color) in enumerate(cards):
        x = left + idx * (w + gap)
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = LIGHT; s.line.color.rgb = BORDER
        cap = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, w, Inches(0.42))
        cap.fill.solid(); cap.fill.fore_color.rgb = color; cap.line.fill.background()
        add_text(slide, x + Inches(0.12), top + Inches(0.08), w - Inches(0.2), Inches(0.22), head, 16, LIGHT, True)
        add_text(slide, x + Inches(0.14), top + Inches(0.52), w - Inches(0.25), h - Inches(0.6), '\n'.join('• ' + l for l in lines), 13, RGBColor(45,45,52))
    add_footer(slide)


def add_table_like(slide, left, top, widths, rows, fills=None, font_size=12):
    x = left
    row_h = Inches(0.52)
    for r_idx, row in enumerate(rows):
        x = left
        for c_idx, txt in enumerate(row):
            w = widths[c_idx]
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top + row_h * r_idx, w, row_h)
            s.fill.solid()
            fc = LIGHT
            if fills and fills[r_idx][c_idx]:
                fc = fills[r_idx][c_idx]
            elif r_idx == 0:
                fc = NAVY
            elif c_idx == 0:
                fc = RGBColor(239, 243, 252)
            s.fill.fore_color.rgb = fc
            s.line.color.rgb = BORDER
            tb = slide.shapes.add_textbox(x + Inches(0.05), top + row_h * r_idx + Inches(0.05), w - Inches(0.1), row_h - Inches(0.08))
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(font_size)
            r.font.bold = (r_idx == 0 or c_idx == 0)
            r.font.color.rgb = LIGHT if r_idx == 0 else RGBColor(38, 43, 51)
            x += w

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_cover(slide)
add_footer(slide, '韩国TikTok Shop从0到1指南 v2.0 · 修订版封面')

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '1. 路径分流页：先别混，先选路', '核心修正页')
add_flow_cards(slide, '1. 路径分流页：先别混，先选路', [
    ('A. 韩国本土站（待官方进一步明确）', ['适合：有韩国本地公司、仓、客服、税务能力', '现状：信息披露不稳定，需以 TikTok 官方公告/邀请为准', '判断：不是新手第一条路'], PURPLE),
    ('B. 韩国主体做 SEA 跨境（主路径）', ['适合：想快速起量、先做东南亚', '优势：韩国主体可对接东南亚市场，更可执行', '判断：当前最值得从 0 到 1 跑通'], TEAL),
    ('C. 反向海淘（韩国→中国）', ['适合：有韩国商品资源/内容种草能力', '打法：小红书/抖音种草 + 预售/代发', '判断：不是 TikTok Shop 主线，但能补现金流'], ORANGE)
], top=Inches(1.2), h=Inches(2.0))
add_bullet_box(slide, Inches(0.7), Inches(3.55), Inches(5.95), Inches(2.6), '新手决策树（身份→预算→品类→市场）', [
    '身份：若无可用韩国主体/长期签证合规条件，先不要碰韩国本土闭环。',
    '预算：<300万韩元先轻SKU + 小批量测品；>800万韩元再考虑韩国本地仓。',
    '品类：优先低客诉、低监管、低侵权；避开功能性化妆品、食品、儿童高监管品。',
    '市场：先 SEA 跨境跑模型，再评估韩国本土闭环或反向海淘。'
], accent=BLUE, font_size=14)
add_bullet_box(slide, Inches(6.9), Inches(3.55), Inches(5.7), Inches(2.6), '一句话结论', [
    '把“韩国本地机会”当中长期 option，把“韩国主体做 SEA”当当前主战场。',
    '如果你连选品、内容、履约都没跑通，先别上高监管本地站，容易烧钱。',
    '需要练手平台时，可先用 Amazon 低难度站点做基本功。'
], accent=RED, font_size=14)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '2. 韩国 TikTok 市场格局：有机会，但别把它想得太成熟')
add_bullet_box(slide, Inches(0.6), Inches(1.2), Inches(6.0), Inches(2.3), '市场判断', [
    '韩国用户对短视频内容消费强，但电商心智长期被 Coupang、Naver、11Street 等平台教育。',
    'TikTok 的优势是“内容驱动冲动购买”，不是“搜索驱动的标准化货架”。',
    '因此：爆品机会存在，但转化效率、履约体验、平台规则波动都比成熟货架平台更敏感。'
], accent=BLUE, font_size=14)
add_bullet_box(slide, Inches(6.75), Inches(1.2), Inches(6.0), Inches(2.3), '风险提示', [
    '不要把“韩国用户喜欢 TikTok”直接等同于“韩国本土 TikTok Shop 已完全成熟”。',
    '支付、物流、税务、消费者保护、平台灰度开放节奏，都可能影响落地细节。',
    '凡是涉及韩国本土闭环能力，必须按最新官方 Seller Center 页面/邀请政策复核。'
], accent=ORANGE, font_size=14)
add_bullet_box(slide, Inches(0.6), Inches(3.8), Inches(12.15), Inches(2.1), '操盘建议', [
    '韩国市场更适合：高颜值家居、生活方式、轻小件配饰、宠物用品、礼品类、非强功效型美妆工具。',
    '不适合新手：功能性化妆品、食品保健、儿童用品、带认证的电子产品、任何灰色 IP 周边。',
    '先用“内容测试→小单履约→复购验证”跑通，再谈扩仓、扩人、扩国家。'
], accent=TEAL, font_size=15)
add_footer(slide)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '3. Korea → SEA Cross-border：当前最可执行的主路径')
add_bullet_box(slide, Inches(0.7), Inches(1.25), Inches(5.75), Inches(2.4), '为什么它是主路径', [
    '相比韩国本土闭环，SEA 市场对“内容电商 + 跨境履约”接受度更高。',
    '韩国主体切入，能借“韩系审美/韩货心智/韩国公司背书”做差异化。',
    '新手最重要的是跑通模型，不是第一天就挑战最复杂监管。'
], accent=TEAL, font_size=14)
add_bullet_box(slide, Inches(6.7), Inches(1.25), Inches(6.0), Inches(2.4), '推荐打法', [
    '先选 1–2 个国家做 MVP，不建议一开始同时开太多市场。',
    '优先轻小件、低售后、内容好演示的 SKU：家居小工具、收纳、宠物、小配件。',
    '使用统一素材底稿，再按越南/泰国/菲律宾语言与价格带做二次本地化。'
], accent=BLUE, font_size=14)
add_bullet_box(slide, Inches(0.7), Inches(3.95), Inches(12.0), Inches(2.0), '底线提醒', [
    'SEA 跨境“可执行”不等于“无门槛”：仍然要验证站点开放资格、跨境收款、类目限制、当地税务与进口要求。',
    '如果官方路径更新，优先看 TikTok 官方入驻页和 Seller Academy，而不是中介口径。'
], accent=RED, font_size=15)
add_footer(slide)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '4. 费用结构：先算毛利，再谈起盘')
add_table_like(slide, Inches(0.7), Inches(1.25), [Inches(2.4), Inches(3.2), Inches(3.0), Inches(3.25)], [
    ['成本项', '常见区间/判断', '关键备注', '新手建议'],
    ['平台佣金', '按站点/类目浮动，常见约 2%–15%', '以官方费率表为准', '测品时按高值预估，别按最低值算'],
    ['支付/收款', '支付处理+提现+汇损合计常见 2%–4%', '不同 PSP / 币种差异大', '必须把汇率损耗单独列'],
    ['物流', '跨境头程 + 尾程 + 逆向退货', '低客单价最容易被物流吃掉', '轻小件优先'],
    ['内容投放', '短视频样品 / 达人佣金 / 广告加热', '经常被新手低估', '至少预留首月试错预算'],
    ['仓储', '韩国本地仓更快但更贵；跨境仓更灵活但时效弱', '按周转率算，不只看单价', '先小量试仓']
], font_size=12)
add_bullet_box(slide, Inches(0.72), Inches(4.85), Inches(12.0), Inches(1.2), '免责声明（必须强化）', [
    'P5 已删除“必须接入 Naver Pay / KakaoPay”的错误表述。它们很重要，但是否必须接入，取决于站点开放能力、服务商和当期平台配置，不可绝对化。'
], accent=ORANGE, font_size=13)
add_footer(slide)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '5. 新手适合品类 vs 不适合品类（新增）')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(4.8), '适合新手切入', [
    '家居收纳 / 厨房小工具：内容展示直观，决策快。',
    '手机壳、桌面配件、线材收纳：轻、小、易寄。',
    '宠物玩具/宠物清洁小用品：复购潜力强。',
    '礼品、小饰品、节日氛围商品：容易做情绪价值内容。',
    '非功能性美妆工具：如化妆刷收纳、美妆镜、发饰。'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(6.75), Inches(1.2), Inches(6.0), Inches(4.8), '不适合新手切入', [
    '功能性化妆品：涉及 MFDS/宣称/成分风险，新手不要从这里开局。',
    '食品、保健品、儿童用品：高投诉 + 高监管。',
    '带电/无线/认证类电子：认证、售后、运输都麻烦。',
    '服装鞋靴：尺码退货和库存压力大。',
    '非授权 K-Pop 周边 / 擦边 IP：侵权风险高，已明确删除。'
], accent=RED, font_size=15)
add_footer(slide)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '6. 韩国企业注册流程：先把主体搭稳')
add_bullet_box(slide, Inches(0.8), Inches(1.2), Inches(3.8), Inches(4.9), '流程顺序', [
    '1) 确认签证/身份允许创业',
    '2) 设立个人事业者或法人',
    '3) 开企业银行账户',
    '4) 准备税务登记/发票/合同模板',
    '5) 准备 TikTok / 收款 / 仓配资料'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(4.85), Inches(1.2), Inches(3.8), Inches(4.9), '常见材料', [
    '营业执照/法人资料',
    '韩国手机号与联系邮箱',
    '银行账户与可扣费信用卡',
    '受益所有人/实际经营人信息',
    '办公地址、退货地址、经营范围'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(8.9), Inches(1.2), Inches(3.8), Inches(4.9), '常见踩坑', [
    '以不具备经营资格的签证身份直接开店。',
    '主体、收款、仓库、联系人信息不一致。',
    '把个人账户拿来收商业流水。',
    '没把税务代理和后续申报安排好就先开卖。'
], accent=ORANGE, font_size=15)
add_footer(slide)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '7. TikTok Shop 入驻（A）：SEA 跨境版，可执行')
add_bullet_box(slide, Inches(0.75), Inches(1.2), Inches(5.8), Inches(4.9), '建议流程', [
    '确认韩国主体是否满足当前跨境站点开放条件。',
    '准备企业资料、联系人、收款账户、退货地址、主营类目说明。',
    '先开 1–2 个目标市场，优先菲律宾/泰国等更适合练手的站点。',
    '先上 10–20 个标准化 SKU，配合 20–30 条测试素材跑点击与转化。',
    '把客服 SLA、逆向退货和差评 SOP 同时搭起来。'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(4.9), '审核重点', [
    '主体真实性、类目适配、素材合规、物流承诺和收款链路。',
    '高风险词：功效夸大、医疗暗示、侵权素材、灰色品牌擦边。',
    '如果走服务商/代运营路径，必须分清：谁是店铺主体，谁负责税务，谁承担违规。'
], accent=BLUE, font_size=15)
add_footer(slide)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '8. TikTok Shop 入驻（B）：韩国本土版，暂按“待官方定型”处理')
add_bullet_box(slide, Inches(0.8), Inches(1.2), Inches(5.7), Inches(4.9), '为什么要单独拆页', [
    'v1.2 最大问题就是把“韩国本土站”与“韩国主体做 SEA”混为一谈。',
    '韩国本土闭环涉及更强的本地支付、物流、客服、消费者保护与税务能力。',
    '而且平台开放节奏、具体规则披露，可能阶段性变化。'
], accent=PURPLE, font_size=15)
add_bullet_box(slide, Inches(6.75), Inches(1.2), Inches(5.7), Inches(4.9), '当前正确说法', [
    '可以持续关注，但不要把它当成当前新手的默认入口。',
    '任何“已经完全开放、流程固定、某支付方式必须接入”的说法，都要重新核验。',
    '若你已经有韩国本地仓、韩语客服、税务与退货能力，再考虑试点。'
], accent=RED, font_size=15)
add_footer(slide)

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '9. 仓储与物流：速度是转化率的一部分')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(4.0), Inches(4.9), '方案 A：韩国本地仓', [
    '优点：时效快、退货处理顺、用户体验更稳定。',
    '缺点：压货、仓租、尾程、库存不准都会吃利润。',
    '适合：爆款验证后，稳定周销再上。'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(4.95), Inches(1.2), Inches(4.0), Inches(4.9), '方案 B：跨境直发 / 中转仓', [
    '优点：试错便宜，适合前期测品。',
    '缺点：时效慢、退货难、体验差。',
    '适合：轻小件、小批量、多款测试。'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(9.2), Inches(1.2), Inches(3.5), Inches(4.9), '物流 KPI', [
    '准时发货率',
    '签收时效',
    '破损率/丢件率',
    '逆向退货成本',
    '因物流导致的差评率'
], accent=ORANGE, font_size=15)
add_footer(slide)

# Slide 11
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '10. 内容创作：不是把中文视频翻成韩文就行')
add_bullet_box(slide, Inches(0.7), Inches(1.25), Inches(5.9), Inches(4.8), '本地化细节', [
    '前 3 秒必须更快：韩国用户刷得快，开场要直接给结果、反差或痛点。',
    '镜头语言要更干净：少土味字幕，多生活方式表达。',
    '评论区脚本要提前写：常见质疑、发货时效、退换政策必须有人设化回复。',
    '同一条素材至少做 3 个版本：封面、开头、价格锚点不同。'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(6.85), Inches(1.25), Inches(5.75), Inches(4.8), '内容选题模板', [
    'Before / After（使用前后）',
    '痛点实测（真的省时间吗？）',
    '韩国生活方式场景化（宿舍/厨房/办公桌）',
    '礼物/节日/情侣/宠物等高分享场景',
    'UGC 口碑二创，而不是纯硬广堆砌'
], accent=TEAL, font_size=15)
add_footer(slide)

# Slide 12
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '11. 客服与售后：差评不是客服问题，是系统问题')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(5.85), Inches(4.9), '客服体系', [
    '设定 24 小时内响应目标；高峰期要有 FAQ + AI + 人工兜底。',
    '把发货、延误、退款、换货、质保话术标准化。',
    '敏感投诉必须升级：物流破损、功效争议、侵权投诉、平台处罚。'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(6.8), Inches(1.2), Inches(5.85), Inches(4.9), '售后底线', [
    '先定义“哪些商品支持无理由、哪些只能质量退”。',
    '发货前拍照/录视频留档，给恶意退货留证据。',
    '把差评原因归因到 SKU、素材、物流、客服四类，而不是只怪客服。'
], accent=RED, font_size=15)
add_footer(slide)

# Slide 13
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '12. 财务合规与 VAT：这页必须比 v1.2 更谨慎')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(12.0), Inches(2.1), '核心免责声明', [
    'VAT / GST / 销售税规则会因主体所在地、仓储地、销售地、平台代扣模式、进口申报方式而变化。下面仅作经营判断框架，不构成税务意见；正式开卖前，必须让韩国税务师 + 目标市场本地顾问复核。'
], accent=ORANGE, font_size=15)
add_bullet_box(slide, Inches(0.7), Inches(3.55), Inches(5.8), Inches(2.3), '实操重点', [
    '主体、收款、开票、报税口径要统一。',
    '跨境平台代扣 ≠ 你可以完全不管税务台账。',
    '每周看毛利，不要等季度报税才发现亏损。'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(6.8), Inches(3.55), Inches(5.9), Inches(2.3), '红线提醒', [
    '个人账户收商业款、混用公私流水、虚报货值、拿“代扣税”当全部税务答案，都是高危操作。',
    '英国/美国亚马逊还涉及 VAT / Sales Tax、EPR、进口商记录等额外问题。'
], accent=RED, font_size=15)
add_footer(slide)

# Slide 14
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '13. 新手冷启动 Q&A：先求跑通，不求一步到位')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(3.9), Inches(4.9), 'Q1：0 单正常吗？', [
    '正常。新号需要时间建立标签。',
    '先看点击率、完播率、加购率，不要只看出单。',
    '同款 3–5 个素材版本轮测。'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(4.85), Inches(1.2), Inches(3.9), Inches(4.9), 'Q2：要不要先投广告？', [
    '可以，但只给“已经有自然点击苗头”的素材加热。',
    '不要拿广告去救一个没有内容吸引力的产品。',
    '先设止损线。'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(9.0), Inches(1.2), Inches(3.7), Inches(4.9), 'Q3：多久扩品/扩站？', [
    '至少先跑出 1 个可复制 SKU。',
    '单店 SOP 稳定后再扩国家。',
    '没有后台数据就扩张，通常死得更快。'
], accent=ORANGE, font_size=15)
add_footer(slide)

# Slide 15
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '14. 恶意退货 / 恶意投诉应对')
add_bullet_box(slide, Inches(0.7), Inches(1.25), Inches(5.85), Inches(4.8), '恶意退货 SOP', [
    '高风险商品打包全程留档。',
    '收回退货时开箱录像、称重、拍细节。',
    '把“可二次销售 / 不可二次销售”标准写清。',
    '证据链齐全后通过平台工单申诉。'
], accent=RED, font_size=15)
add_bullet_box(slide, Inches(6.8), Inches(1.25), Inches(5.85), Inches(4.8), '恶意投诉 SOP', [
    '先分类型：侵权、虚假宣传、竞品举报、物流时效。',
    '侵权类先下架核查，不要嘴硬。',
    '能提供原创素材、授权书、供货链证明的，胜率更高。',
    '维权记录要沉淀成内部黑名单与素材库。'
], accent=BLUE, font_size=15)
add_footer(slide)

# Slide 16
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '15. 华人身份合规：F-5 是重点，但不是唯一答案')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(4.0), Inches(4.9), '高风险身份', [
    'D-2 / D-4 / D-10 等通常不适合直接从事盈利经营。',
    '不要用侥幸心理拿别人主体、自己暗中经营。',
    '签证问题比店铺问题更严重。'
], accent=RED, font_size=15)
add_bullet_box(slide, Inches(4.95), Inches(1.2), Inches(4.0), Inches(4.9), '相对稳妥身份', [
    'F-5（永驻）通常具备较强经营稳定性。',
    'F-6、部分 F-2 类型通常更适合创业，但仍需核对个人条件。',
    '最稳的做法：先确认身份权限，再注册主体。'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(9.2), Inches(1.2), Inches(3.5), Inches(4.9), '一句提醒', [
    '主体合规 > 选品速度。',
    '别为了快，把签证和税务埋雷。'
], accent=ORANGE, font_size=15)
add_footer(slide)

# Slide 17
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '16. 华人双向机遇：韩国→SEA / 反向海淘')
add_bullet_box(slide, Inches(0.7), Inches(1.25), Inches(5.85), Inches(4.8), '方向 1：韩国→SEA', [
    '卖“韩系审美 + 韩国主体背书 + 内容电商表达”。',
    '适合做家居、配件、宠物、小礼品、非功效型美妆工具。',
    '重点是多语言内容和跨境履约，而不是压大库存。'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(6.8), Inches(1.25), Inches(5.85), Inches(4.8), '方向 2：韩国→中国反向海淘', [
    '适合小众品牌、线下发现、买手式内容。',
    '玩法是内容种草 + 预售/代发，不一定依赖 TikTok Shop。',
    '优点是更懂中韩语言与文化落差，容易做出差异化内容。'
], accent=BLUE, font_size=15)
add_footer(slide)

# Slide 18
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '17. 建议时间表：别把 2 个月的事想成 2 周')
add_table_like(slide, Inches(0.7), Inches(1.25), [Inches(2.2), Inches(4.0), Inches(3.1), Inches(2.9)], [
    ['阶段', '目标', '关键动作', '通过标准'],
    ['第 1–2 周', '确定路径', '身份核查、主体准备、市场选择、预算表', '知道先做哪条路'],
    ['第 3–6 周', '完成准备', '开主体/收款/物流方案/素材规划', '能提交入驻'],
    ['第 2–3 月', '冷启动', '上 SKU、拍素材、测广告、跑客服 SOP', '出现首批稳定订单'],
    ['第 3–4 月', '优化模型', '砍差品、补高转化品、优化履约', '毛利结构变清晰'],
    ['第 4 月以后', '放大复制', '扩站点、扩达人、扩库存', '在已有利润上扩张']
], font_size=12)
add_footer(slide)

# Slide 19
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '18. 开业核对清单：分层做，不要一起炸')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(4.0), Inches(5.0), '开店前', [
    '身份/签证是否允许经营',
    '主体、银行、收款是否打通',
    '目标市场与目标品类是否选定',
    '税务代理 / 记账方案是否确定'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(4.95), Inches(1.2), Inches(4.0), Inches(5.0), '出单前', [
    '至少 10–20 个可卖 SKU',
    '至少 20–30 条素材底稿',
    '客服 FAQ / 退换货规则 / 差评 SOP',
    '发货时效与退货地址写清'
], accent=TEAL, font_size=15)
add_bullet_box(slide, Inches(9.2), Inches(1.2), Inches(3.5), Inches(5.0), '放大前', [
    '已有盈利 SKU',
    '退款率/投诉率可控',
    '广告投产比可验证',
    '现金流能覆盖补货周期'
], accent=ORANGE, font_size=15)
add_footer(slide)

# Slide 20 Amazon comparison
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '19. 亚马逊 5 国机会对比（新增）：把“市场”与“卖家来源地”分清')
rows = [
    ['国家', '市场成熟度', '入驻难度', '费用/模式', '热门方向', '对华人卖家友好度'],
    ['菲律宾', '无本地 Amazon 站；更像卖家来源地', '个人中低 / 企业中', '做 US/UK 站为主；月费按目标站点', '家居、轻小件、OEM', '高'],
    ['泰国', '无本地站；供应链与工厂资源好', '个人中 / 企业中', '适合供货到 US/UK/JP', '家居、汽配、宠物、手工类', '中高'],
    ['越南', '无本地站；制造增长快', '个人中 / 企业中', '适合跨境卖全球', '家具、家居、纺织、户外', '中高'],
    ['英国', '成熟站点', '个人中 / 企业中高', '£0.75/件 或 £25/月+VAT；佣金多为 8%–15%', '家居、低价日用品、配件', '中'],
    ['美国', '最成熟、最大盘', '个人中 / 企业中高', '$0.99/件 或 $39.99/月；佣金多为 8%–15%', '家居、配件、宠物、工具', '中']
]
add_table_like(slide, Inches(0.35), Inches(1.2), [Inches(1.1), Inches(2.1), Inches(1.55), Inches(2.8), Inches(2.35), Inches(2.75)], rows, font_size=11)
add_bullet_box(slide, Inches(0.6), Inches(4.85), Inches(12.1), Inches(1.25), '关键纠偏', [
    '越南 / 泰国 / 菲律宾目前并不是“本地亚马逊消费大站”，而是更适合作为卖家来源地、供应链来源地或低成本练手起点。普通第三方卖家普遍不适用所谓 Refrigerated Storage；冷藏/冷冻多属特殊项目，不是新手常规路径。'
], accent=RED, font_size=13)
add_footer(slide, '费用参考：Amazon Sell US/UK 公共 pricing 页面；东南亚三国按“卖家来源地做全球站”口径整理。')

# Slide 21
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '20. 推荐练手顺序：菲律宾 → 泰国 → 越南 → 英国 → 美国')
chart_data = CategoryChartData()
chart_data.categories = ['菲律宾', '泰国', '越南', '英国', '美国']
chart_data.add_series('推荐指数（越高越适合先练手）', (90, 82, 75, 58, 50))
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(1.35), Inches(6.0), Inches(4.8), chart_data).chart
chart.has_legend = False
chart.value_axis.maximum_scale = 100
chart.value_axis.minimum_scale = 0
chart.category_axis.tick_labels.font.size = Pt(12)
chart.value_axis.tick_labels.font.size = Pt(10)
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
plot.series[0].format.fill.solid()
plot.series[0].format.fill.fore_color.rgb = BLUE
add_bullet_box(slide, Inches(7.1), Inches(1.35), Inches(5.4), Inches(4.8), '排序逻辑', [
    '菲律宾：英语环境相对友好、练手心理门槛低。',
    '泰国：供应链和跨境经验较多，竞争比美国轻。',
    '越南：制造潜力强，但执行细节更吃本地资源。',
    '英国：成熟但 VAT/EPR/欧洲合规复杂度更高。',
    '美国：最大盘，也最卷、最贵、最需要系统能力。'
], accent=ORANGE, font_size=15)
add_footer(slide)

# Slide 22
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, '21. 学习资源 + 22. 结束页（行动清单）')
add_bullet_box(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(4.9), '建议盯住的官方资源', [
    'TikTok Shop Seller Center / Seller Academy：入驻、费率、政策更新。',
    'Amazon Sell US / UK Pricing：月费、佣金、FBA费用。',
    '韩国税务师/会计事务所：主体、VAT、跨境流水设计。',
    'MFDS / 韩国国税厅 / 各市场海关与税务官网：高监管品类必须查原文。'
], accent=BLUE, font_size=15)
add_bullet_box(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(4.9), '现在就执行的 5 步', [
    '1) 先选路径：韩国本地 / SEA 跨境 / 反向海淘。',
    '2) 先砍掉高监管品类，留下 3 个可测方向。',
    '3) 用预算表算清：佣金、物流、汇损、退货。',
    '4) 先做最小可行店铺，而不是豪华版店铺。',
    '5) 若需要练手，优先从菲律宾 → 泰国 → 越南的亚马逊全球卖家路径开始。'
], accent=TEAL, font_size=15)
add_text(slide, Inches(0.9), Inches(6.35), Inches(11.8), Inches(0.4), '一句结论：韩国本土机会可以看，但当前最现实的 0→1 路线，依然是“韩国主体 + SEA 跨境 + 低监管品类 + 小步快跑”。', 18, NAVY, True, PP_ALIGN.CENTER)
add_footer(slide, 'End · v2.0')

prs.save(OUTPUT)
print(OUTPUT)
