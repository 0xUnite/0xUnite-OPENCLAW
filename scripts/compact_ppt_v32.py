from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pathlib import Path

SRC = Path('/Users/sudi/.openclaw/workspace/韩国TikTokShop从0到1指南_最终v3版.pptx')
DST = Path('/Users/sudi/.openclaw/workspace/韩国TikTokShop从0到1指南_v3.2.pptx')

prs = Presentation(str(SRC))

EMU = 914400

def set_text_style(shape, size=None, bold=None, align=None, font_name='PingFang SC'):
    if not getattr(shape, 'has_text_frame', False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    # tighter default margins
    narrow = shape.width < Inches(2.2) or shape.height < Inches(0.5)
    tf.margin_left = Inches(0.05 if not narrow else 0.035)
    tf.margin_right = Inches(0.05 if not narrow else 0.035)
    tf.margin_top = Inches(0.03 if not narrow else 0.02)
    tf.margin_bottom = Inches(0.03 if not narrow else 0.02)
    for p in tf.paragraphs:
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        p.line_spacing = 1.05
        if align is not None:
            p.alignment = align
        for r in p.runs:
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if font_name:
                r.font.name = font_name
        if not p.runs and size is not None:
            p.font.size = Pt(size)
            p.font.bold = bool(bold)
            p.font.name = font_name


def move(shape, dx=0, dy=0, dw=0, dh=0):
    shape.left = max(0, shape.left + int(dx))
    shape.top = max(0, shape.top + int(dy))
    shape.width = max(int(Inches(0.08)), shape.width + int(dw))
    shape.height = max(int(Inches(0.08)), shape.height + int(dh))


def compact_standard_slide(slide):
    for idx, shape in enumerate(slide.shapes, start=1):
        if not getattr(shape, 'has_text_frame', False):
            # compact large content cards a bit upward/downward
            if shape.top > Inches(0.95) and shape.top < Inches(6.0):
                if shape.height > Inches(3.8):
                    move(shape, dy=-Inches(0.18), dh=Inches(0.30))
                elif shape.height > Inches(0.45):
                    move(shape, dy=-Inches(0.14), dh=Inches(0.06))
            elif shape.top >= Inches(6.0) and shape.top < Inches(6.85):
                move(shape, dy=-Inches(0.12), dh=Inches(0.04))
            continue

        text = shape.text.strip()
        top = shape.top / EMU
        left = shape.left / EMU
        height = shape.height / EMU
        width = shape.width / EMU

        # global text compaction
        set_text_style(shape, size=10.5, bold=False)

        # reduce whitespace by moving content up slightly
        if top > 0.95 and top < 6.0:
            if height > 3.0:
                move(shape, dy=-Inches(0.18), dh=Inches(0.32))
            elif height > 0.42:
                move(shape, dy=-Inches(0.14), dh=Inches(0.08))
        elif top >= 6.0 and top < 6.85:
            move(shape, dy=-Inches(0.12), dh=Inches(0.06))

        # role-based styling
        if idx == 2 or (top < 0.42 and width > 5.5):
            set_text_style(shape, size=20, bold=True)
            shape.height = max(shape.height, int(Inches(0.34)))
        elif idx == 3 or (top < 0.8 and height < 0.28 and left < 8.5):
            set_text_style(shape, size=10.5, bold=False)
        elif idx == 5 or (left > 10.5 and top < 0.45):
            set_text_style(shape, size=10, bold=True, align=PP_ALIGN.CENTER)
        elif top > 6.85:
            set_text_style(shape, size=8.2, bold=False)
        elif top >= 6.1 and height <= 0.5:
            # bottom key point / label
            if '核心要点' in text or '本页要点' in text:
                set_text_style(shape, size=9.4, bold=True)
                shape.height = max(shape.height, int(Inches(0.34)))
            else:
                set_text_style(shape, size=9.8, bold=True)
                shape.height = max(shape.height, int(Inches(0.22)))
        elif top >= 5.8:
            set_text_style(shape, size=10.2, bold=False)
            shape.height = max(shape.height, int(Inches(0.24)))
        elif 1.0 <= top <= 1.45 and height <= 0.42:
            set_text_style(shape, size=11.4, bold=True)
        elif width < 2.2 and height < 0.8:
            set_text_style(shape, size=9.7, bold=True)
        else:
            set_text_style(shape, size=10.4, bold=False)


def compact_cover(slide):
    # Update version strings
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            if 'v3.1 审核增强版' in shape.text:
                shape.text = shape.text.replace('v3.1 审核增强版', 'v3.2 紧凑版')
            if '最终审核增强版 v3.1' in shape.text:
                shape.text = shape.text.replace('最终审核增强版 v3.1', '最终紧凑版 v3.2')

    for idx, shape in enumerate(slide.shapes, start=1):
        if getattr(shape, 'has_text_frame', False):
            set_text_style(shape, size=10.5, bold=False)

    # tighter layout on cover
    for shape in slide.shapes:
        if shape.top > Inches(0.8) and shape.top < Inches(5.8):
            if shape.height > Inches(1.0):
                move(shape, dy=-Inches(0.16), dh=Inches(0.24))
            else:
                move(shape, dy=-Inches(0.10), dh=Inches(0.04))

    # specific cover typography
    if len(slide.shapes) >= 22:
        set_text_style(slide.shapes[3], size=9.5, bold=True, align=PP_ALIGN.CENTER)
        set_text_style(slide.shapes[4], size=22, bold=True)
        set_text_style(slide.shapes[5], size=11, bold=False)
        set_text_style(slide.shapes[7], size=10.5, bold=True)
        set_text_style(slide.shapes[8], size=10.5, bold=False)
        set_text_style(slide.shapes[9], size=10.0, bold=True)
        set_text_style(slide.shapes[11], size=9.5, bold=True, align=PP_ALIGN.CENTER)
        set_text_style(slide.shapes[13], size=9.5, bold=True, align=PP_ALIGN.CENTER)
        set_text_style(slide.shapes[15], size=9.5, bold=True, align=PP_ALIGN.CENTER)
        set_text_style(slide.shapes[17], size=10.0, bold=True)
        set_text_style(slide.shapes[18], size=10.3, bold=False)
        set_text_style(slide.shapes[19], size=8.2, bold=False, align=PP_ALIGN.RIGHT)
        set_text_style(slide.shapes[20], size=9.0, bold=True, align=PP_ALIGN.CENTER)
        set_text_style(slide.shapes[21], size=9.6, bold=True)

compact_cover(prs.slides[0])
for i in range(1, len(prs.slides)):
    compact_standard_slide(prs.slides[i])

# version/footer refresh without changing substantive content
for i in range(1, len(prs.slides)):
    slide = prs.slides[i]
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and 'v3.1 审核增强版' in shape.text:
            shape.text = shape.text.replace('v3.1 审核增强版', 'v3.2 紧凑版')
            set_text_style(shape, size=8.2, bold=False)

prs.save(str(DST))
print(DST)
