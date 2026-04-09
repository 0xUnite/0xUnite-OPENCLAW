from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

OUTPUT = '/Users/sudi/.openclaw/workspace/韩国TikTokShop从0到1指南_v3.0.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT = 'PingFang SC'
FONT_EN = 'Aptos'

# Color system
BG = RGBColor(0xF7, 0xF8, 0xFA)
COVER_BG = RGBColor(0xF5, 0xF7, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xE6, 0xEA, 0xF0)
TITLE = RGBColor(0x1F, 0x29, 0x37)
SUBTITLE = RGBColor(0x33, 0x41, 0x55)
BODY = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x64, 0x74, 0x8B)
BLUE = RGBColor(0x5B, 0x8D, 0xEF)
BLUE_BG = RGBColor(0xDC, 0xEB, 0xFF)
BLUE_TXT = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x2E, 0xC5, 0xA7)
GREEN_BG = RGBColor(0xE8, 0xFB, 0xF4)
GREEN_TXT = RGBColor(0x0F, 0x9F, 0x7F)
PINK = RGBColor(0xFF, 0x6B, 0x8A)
PINK_BG = RGBColor(0xFF, 0xE7, 0xEE)
PINK_TXT = RGBColor(0xD9, 0x48, 0x5F)
YELLOW_BG = RGBColor(0xFF, 0xF4, 0xD6)
YELLOW_TXT = RGBColor(0x9A, 0x67, 0x0A)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text='', size=12, color=BODY, bold=False,
                font=FONT, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def add_paragraphs(tb, lines, size=11.5, color=BODY, bullet=True, level=0, space_after=2):
    tf = tb.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = level
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(space_after)
        p.bullet = bullet


def add_header(slide, title, subtitle=None, page=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.82))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CARD
    bar.line.color.rgb = DIVIDER
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(8.6), Inches(0.3), title, size=22, color=TITLE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.49), Inches(6.5), Inches(0.2), subtitle, size=11.5, color=MUTED)
    if page is not None:
        add_tag(slide, Inches(11.9), Inches(0.18), Inches(0.9), Inches(0.34), f'{page:02d}', BLUE_BG, BLUE_TXT, 11.5)


def add_footer(slide, text='韩国TikTok Shop从0到1指南 · v3.0 全新配色版 · 仅供经营判断参考，请以平台/税务官方规则为准'):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.01))
    line.fill.solid(); line.fill.fore_color.rgb = DIVIDER; line.line.fill.background()
    add_textbox(slide, Inches(0.52), Inches(6.98), Inches(12.0), Inches(0.18), text, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)


def add_tag(slide, left, top, width, height, text, bg_color, text_color, size=11.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    add_textbox(slide, left, top + Inches(0.01), width, height - Inches(0.02), text, size=size, color=text_color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    return shape


def add_card(slide, left, top, width, height, title, body_lines=None, accent=BLUE, title_size=16, body_size=11.5, fill=CARD):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = DIVIDER
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.32), Inches(0.35), title, size=title_size, color=SUBTITLE, bold=True)
    if body_lines:
        tb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.52), width - Inches(0.32), height - Inches(0.62))
        add_paragraphs(tb, body_lines, size=body_size, color=BODY, bullet=True, space_after=2)
    return card


def add_note_box(slide, left, top, width, height, title, lines, bg_color, title_color, body_color=BODY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = bg_color
    box.line.fill.background()
    add_textbox(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.3), Inches(0.28), title, size=13, color=title_color, bold=True)
    tb = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.42), width - Inches(0.32), height - Inches(0.52))
    add_paragraphs(tb, lines, size=11.2, color=body_color, bullet=True, space_after=2)


def add_two_col_slide(title, subtitle, left_title, left_lines, right_title, right_lines, page,
                      left_accent=BLUE, right_accent=GREEN, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, title, subtitle, page)
    add_card(slide, Inches(0.58), Inches(1.15), Inches(6.0), Inches(5.15), left_title, left_lines, accent=left_accent)
    add_card(slide, Inches(6.76), Inches(1.15), Inches(6.0), Inches(5.15), right_title, right_lines, accent=right_accent)
    if note:
        note_title, note_lines, style = note
        bg, color = style
        add_note_box(slide, Inches(0.58), Inches(6.38), Inches(12.18), Inches(0.38), note_title, note_lines, bg, color)
    add_footer(slide)


def add_three_col_slide(title, subtitle, cols, page, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_header(slide, title, subtitle, page)
    x_positions = [Inches(0.5), Inches(4.48), Inches(8.46)]
    for x, (ctitle, lines, accent) in zip(x_positions, cols):
        add_card(slide, x, Inches(1.18), Inches(3.78), Inches(5.05), ctitle, lines, accent=accent, title_size=15)
    if note:
        note_title, note_lines, style = note
        add_note_box(slide, Inches(0.5), Inches(6.34), Inches(11.74), Inches(0.4), note_title, note_lines, style[0], style[1])
    add_footer(slide)


def add_matrix_table(slide, left, top, col_widths, rows, row_h=0.62, font_size=11.2):
    for r_idx, row in enumerate(rows):
        x = left
        for c_idx, cell in enumerate(row):
            w = col_widths[c_idx]
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if r_idx == 0 else MSO_SHAPE.RECTANGLE,
                                           x, top + Inches(row_h * r_idx), w, Inches(row_h - 0.02))
            shape.fill.solid()
            if r_idx == 0:
                shape.fill.fore_color.rgb = BLUE_BG
                t_color = BLUE_TXT
            elif c_idx == 0:
                shape.fill.fore_color.rgb = CARD
                t_color = SUBTITLE
            else:
                shape.fill.fore_color.rgb = CARD
                t_color = BODY
            shape.line.color.rgb = DIVIDER
            add_textbox(slide, x + Inches(0.05), top + Inches(row_h * r_idx) + Inches(0.06), w - Inches(0.1), Inches(row_h - 0.12),
                        cell, size=font_size, color=t_color, bold=(r_idx == 0 or c_idx == 0))
            x += w


# 1 Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, COVER_BG)
hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.72), Inches(8.45), Inches(5.65))
hero.fill.solid(); hero.fill.fore_color.rgb = CARD; hero.line.color.rgb = DIVIDER
side = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.34), Inches(0.72), Inches(3.36), Inches(5.65))
side.fill.solid(); side.fill.fore_color.rgb = RGBColor(0xEE, 0xF4, 0xFF); side.line.fill.background()
add_tag(slide, Inches(0.92), Inches(0.95), Inches(1.95), Inches(0.42), 'v3.0 全新配色版', BLUE_BG, BLUE_TXT, 12)
add_textbox(slide, Inches(0.95), Inches(1.58), Inches(7.5), Inches(1.5), '韩国TikTok Shop从0到1指南', size=26, color=TITLE, bold=True)
add_textbox(slide, Inches(0.98), Inches(2.55), Inches(7.4), Inches(0.85), '韩国主体切入 TikTok Shop 与 Amazon 备选路径\n27页重构版：先讲合规，再讲路径，再讲执行。', size=15, color=BODY)
add_note_box(slide, Inches(0.95), Inches(3.65), Inches(7.5), Inches(1.55), '本版重点', [
    '讲解顺序完全重排：把“华人身份合规”提前到第3页。',
    '全部改用浅色背景 + 白色卡片 + 浅底深字标签。',
    '新增 Amazon 补充路径完整模块，避免“一平台思维”。'
], GREEN_BG, GREEN_TXT)
add_textbox(slide, Inches(9.62), Inches(1.2), Inches(2.6), Inches(0.5), '适用对象', size=18, color=SUBTITLE, bold=True)
add_tag(slide, Inches(9.62), Inches(1.82), Inches(2.3), Inches(0.38), '韩国主体 / 华人创业者', BLUE_BG, BLUE_TXT)
add_tag(slide, Inches(9.62), Inches(2.3), Inches(1.75), Inches(0.38), 'TikTok新手', GREEN_BG, GREEN_TXT)
add_tag(slide, Inches(9.62), Inches(2.78), Inches(1.95), Inches(0.38), '跨境电商转型', PINK_BG, PINK_TXT)
add_note_box(slide, Inches(9.55), Inches(3.45), Inches(2.85), Inches(1.7), '一句结论', [
    '先跑通 SEA 跨境。',
    '韩国本土先观望。',
    'Amazon 是补充，不是同步重仓。'
], YELLOW_BG, YELLOW_TXT)
add_textbox(slide, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.22), '韩国TikTok Shop从0到1指南 · 2026版', size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)

# 2-27
add_three_col_slide(
    '路径分流页', '先分清你在做哪条路，不要把三种逻辑混成一锅粥。', [
        ('A. 韩国本土闭环', ['需要韩国主体、本地仓、韩语客服、退货与税务能力。', '目前更适合已有本地经营能力的人，而不是纯新手。', '核心问题不是“能不能开”，而是“能不能稳住履约与合规”。'], PINK),
        ('B. Korea → SEA 主路径', ['更适合从0到1练手。', '可以利用韩系审美、韩国主体背书与内容表达做差异化。', '重点先跑通选品、内容、履约、客服闭环。'], BLUE),
        ('C. 反向海淘 / 买手模式', ['更像内容种草 + 预售/代发。', '不一定依赖 TikTok Shop。', '适合有韩国货源发现能力的人。'], GREEN)
    ], 2, note=('决策提醒', ['新手默认优先走 B，不要一上来挑战最复杂的本地闭环。'], (YELLOW_BG, YELLOW_TXT))
)

add_three_col_slide(
    '华人身份合规（提前！）', '这页必须提前讲：主体合规，比选品速度重要得多。', [
        ('高风险身份', ['D-2 / D-4 / D-10 等通常不适合直接盈利经营。', '不要抱侥幸心理借名义主体、自己暗中实操。', '签证问题一旦出事，比店铺违规更麻烦。'], PINK),
        ('相对稳妥身份', ['F-5 永驻通常经营稳定性更强。', 'F-6、部分 F-2 类型通常更适合创业，但仍要个案确认。', '最稳做法：先核查身份权限，再注册主体。'], GREEN),
        ('底线建议', ['店铺主体、实际经营人、收款账户尽量一致。', '别用个人账户承接商业流水。', '先咨询行政书士/税务师，再谈扩店。'], BLUE)
    ], 3, note=('一句话', ['主体合规 > 店铺速度 > 爆单幻想。'], (PINK_BG, PINK_TXT))
)

add_three_col_slide(
    '韩国企业注册流程', '先把公司与基础账户搭稳，后续平台审核和财务才不会反复出错。', [
        ('流程顺序', ['确认签证/身份允许经营。', '设立个人事业者或法人。', '开企业银行账户。', '准备税务登记、发票与合同模板。', '再去准备平台与仓配资料。'], BLUE),
        ('所需资料', ['营业执照/法人资料。', '韩国手机号、企业邮箱。', '企业银行账户、可扣费信用卡。', '办公地址、退货地址、受益人信息。'], GREEN),
        ('常见踩坑', ['主体、联系人、仓储、收款信息不一致。', '个人账户混收公司流水。', '没安排税务申报就先开卖。', '材料临时拼凑，导致 KYC 卡住。'], PINK)
    ], 4
)

add_two_col_slide(
    '韩国TikTok内容消费环境', '韩国用户爱刷短视频，但不代表 TikTok Shop 本地电商成熟度已经完全到位。',
    '市场特点', [
        'TikTok 强在内容驱动冲动购买，不是搜索驱动货架电商。',
        '韩国电商长期由 Coupang、Naver 等平台教育，用户对履约体验要求高。',
        '因此内容能起量，但物流、退货、客服更容易决定口碑。'
    ],
    '操盘含义', [
        '适合高颜值、直观演示、轻决策商品。',
        '不适合强功效、重认证、重售后类目。',
        '不要把“内容流量”误判成“平台基础设施已经成熟”。'
    ],
    5, note=('风险提示', ['韩国本土规则与开放节奏会波动，任何绝对化说法都要复核官方信息。'], (PINK_BG, PINK_TXT))
)

add_two_col_slide(
    'Korea → SEA主路径', '这是当前最现实的 0→1 路线：先跑模型，而不是先挑战最高复杂度。',
    '为什么优先走这条路', [
        'SEA 对内容电商和跨境履约的接受度更高。',
        '韩国主体可借韩系审美和韩国背书做差异化。',
        '更适合先验证 SKU、素材、广告和客服闭环。'
    ],
    '执行建议', [
        '先选 1–2 个市场，不要同时铺太多国家。',
        '优先轻小件、低客诉、好演示商品。',
        '统一素材底稿，再做越南/泰国/菲律宾二次本地化。'
    ],
    6, note=('底线提醒', ['“可执行”不等于“无门槛”，仍要核查站点开放、物流与税务条件。'], (YELLOW_BG, YELLOW_TXT))
)

add_two_col_slide(
    '新手适合 / 不适合品类', '先把坑绕开，冷启动成功率会高很多。',
    '适合新手', [
        '家居收纳、厨房小工具：展示直观，决策快。',
        '桌面整理、手机配件、线材收纳：轻小件，物流友好。',
        '宠物用品、小礼品、氛围类商品：情绪价值内容好做。',
        '非功能性美妆工具：比直接卖功效型美妆更稳。'
    ],
    '不适合新手', [
        '化妆品、食品、儿童用品：高监管高投诉。',
        '带电产品、认证类电子：认证和售后都麻烦。',
        '易碎高退货类目：售后和损耗吃利润。',
        '侵权商品、擦边IP：一旦踩线，直接送命。'
    ],
    7, left_accent=GREEN, right_accent=PINK
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '费用结构（含SKU毛利测算案例）', '所有新手都容易高估毛利，低估物流、投放、退货与汇损。', 8)
rows = [
    ['成本项', '常见判断', '备注'],
    ['平台佣金', '按站点/类目浮动', '测算时按偏高值估，不要按最低值做梦。'],
    ['支付/收款/汇损', '常见合计 2%–4%', '不同 PSP 和币种差异很大。'],
    ['物流', '头程 + 尾程 + 退货', '低客单价最容易被物流吃掉。'],
    ['内容/达人/广告', '常被新手漏算', '首月必须留试错预算。'],
    ['仓储/包材/损耗', '与周转率强相关', '不是只看仓租单价。']
]
add_matrix_table(slide, Inches(0.55), Inches(1.2), [Inches(2.0), Inches(3.0), Inches(7.3)], rows, row_h=0.63)
add_note_box(slide, Inches(0.58), Inches(5.12), Inches(5.9), Inches(1.3), 'SKU毛利案例', [
    '售价 19.9 美元；进货 4.2；物流 3.8；平台佣金 2.0；广告/达人 2.6；包材与损耗 0.8。',
    '预计到手毛利 ≈ 6.5 美元，毛利率约 32.7%。',
    '一旦退货率抬升或广告失控，这个利润会非常快地蒸发。'
], GREEN_BG, GREEN_TXT)
add_note_box(slide, Inches(6.65), Inches(5.12), Inches(6.1), Inches(1.3), '利润公式', [
    '到手利润 = 售价 - 佣金 - 支付/汇损 - 物流 - 广告/达人 - 包材 - 退货损耗 - 税务成本。',
    '别只看“出单利润”，要看“退货后利润”。'
], YELLOW_BG, YELLOW_TXT)
add_footer(slide)

add_two_col_slide(
    'TikTok Shop入驻 - SEA跨境版', '这是当前更值得新手投入精力的路线。',
    '建议流程', [
        '确认韩国主体是否满足当前跨境站点开放条件。',
        '准备企业资料、联系人、收款账户、退货地址、主营类目说明。',
        '先开 1–2 个目标市场，先上 10–20 个标准化 SKU。',
        '同步搭好客服 SLA、逆向退货与差评 SOP。'
    ],
    '审核重点', [
        '主体真实性、类目适配、素材合规、物流承诺与收款链路。',
        '高风险词：功效夸大、医疗暗示、侵权素材、灰色品牌擦边。',
        '如果走服务商路径，必须分清：谁是主体，谁负责税务，谁承担违规。'
    ],
    9
)

add_two_col_slide(
    'TikTok Shop入驻 - 韩国本土（观望）', '这不是不能看，而是不该当成新手默认入口。',
    '为什么要谨慎', [
        '韩国本土闭环更吃支付、物流、客服、退货与消费者保护能力。',
        '开放节奏与具体规则可能阶段性变化。',
        '只盯“能不能开店”而忽略运营基础，很容易翻车。'
    ],
    '当前正确姿势', [
        '持续关注官方 Seller Center 与公告。',
        '已有韩国本地仓、韩语客服与税务能力的人可以评估试点。',
        '新手先别把大量预算砸在这条路上。'
    ],
    10, left_accent=PINK, right_accent=BLUE, note=('提醒', ['任何“已经完全开放、流程固定、某支付方式必须接入”的说法，都要重新核验。'], (PINK_BG, PINK_TXT))
)

add_three_col_slide(
    '仓储与物流', '速度、破损率、逆向退货处理，都会直接影响内容转化和评分。', [
        ('韩国本地仓', ['优点：时效快、退货处理顺、用户体验更稳。', '缺点：压货、仓租、尾程和库存错误都会吃利润。', '适合：爆品验证后再上。'], BLUE),
        ('跨境直发 / 中转仓', ['优点：试错成本低，适合前期测品。', '缺点：时效慢、退货难、体验差。', '适合：轻小件、小批量、多款测试。'], GREEN),
        ('物流KPI', ['准时发货率。', '签收时效。', '破损/丢件率。', '逆向退货成本。', '物流导致差评率。'], PINK)
    ], 11
)

add_two_col_slide(
    '内容创作', '不是把中文视频翻成韩文就完了，真正难的是“本地化表达”。',
    '内容方法', [
        '前3秒更快给结果、反差或痛点。',
        '镜头语言更干净，少土味字幕，多生活方式感。',
        '同一条素材至少做 3 个版本：封面、开头、价格锚点都改。'
    ],
    '选题模板', [
        'Before / After 使用前后。',
        '痛点实测：真的省时间吗？',
        '宿舍 / 厨房 / 办公桌等韩国生活场景。',
        '礼物、情侣、宠物、节日等高分享场景。'
    ],
    12, note=('评论区', ['评论区脚本必须提前写，发货时效、退款政策、真假质疑都要有人设化回应。'], (GREEN_BG, GREEN_TXT))
)

add_two_col_slide(
    '客服与售后', '差评不是客服单点问题，而是整个系统设计的结果。',
    '客服体系', [
        '设定 24 小时内响应目标，高峰期用 FAQ + AI + 人工兜底。',
        '把发货、延误、退款、换货、质保话术标准化。',
        '敏感投诉必须升级：物流破损、功效争议、侵权投诉、平台处罚。'
    ],
    '售后底线', [
        '先定义哪些商品支持无理由，哪些只能质量退。',
        '发货前拍照/录视频留档，给恶意退货留证据。',
        '把差评归因到 SKU、素材、物流、客服四类，而不是只怪客服。'
    ],
    13, right_accent=PINK
)

add_two_col_slide(
    '财务合规与间接税', '这一页只讲经营判断框架，不替代正式税务意见。',
    '实操重点', [
        '主体、收款、开票、报税口径尽量统一。',
        '平台代扣税款 ≠ 你可以完全不管税务台账。',
        '每周看毛利，不要等季度报税才发现自己一直亏。'
    ],
    '红线提醒', [
        '个人账户收商业款、混用公私流水、虚报货值都属于高危操作。',
        '跨境站点税制会随主体、仓储地、销售地与进口方式变化。',
        '正式开卖前，请让韩国税务师 + 目标市场本地顾问复核。'
    ],
    14, left_accent=BLUE, right_accent=PINK, note=('免责声明', ['平台代扣 ≠ 合规结束，完整台账必须自己留。'], (YELLOW_BG, YELLOW_TXT))
)

add_three_col_slide(
    '开业核对清单', '分层检查，别把所有风险堆到上线当天。', [
        ('开店前', ['身份/签证是否允许经营。', '主体、银行、收款是否打通。', '目标市场与目标品类是否明确。', '税务代理/记账方案是否确定。'], BLUE),
        ('出单前', ['至少 10–20 个可卖 SKU。', '至少 20–30 条素材底稿。', '客服 FAQ / 退换货规则 / 差评 SOP 已成型。', '发货时效和退货地址写清。'], GREEN),
        ('放大前', ['已有盈利 SKU。', '退款率/投诉率可控。', '广告投产比可验证。', '现金流覆盖补货周期。'], PINK)
    ], 15
)

add_two_col_slide(
    '新手MVP起步方案', '不要上来就做“豪华版店铺”，先做能跑通的最小闭环。',
    'MVP配置', [
        '1 个市场 + 3–10 个 SKU。',
        '20 条以内测试素材。',
        '简单客服 SOP + 退货规则 + 预算表。',
        '小批量备货或中转仓，先看真实转化。'
    ],
    '第一阶段目标', [
        '验证点击率、加购率、退款率，而不是只看 GMV。',
        '找到 1 个可复制 SKU。',
        '把素材、广告、客服、物流闭环跑顺。',
        '3 周到 6 周内完成第一轮复盘。'
    ],
    16, note=('一句人话', ['新手第一目标不是赚钱很多，而是别把钱烧在错误路径上。'], (YELLOW_BG, YELLOW_TXT))
)

add_three_col_slide(
    '新手冷启动Q&A', '先求跑通，再求放大。', [
        ('Q1：0单正常吗？', ['正常。新号需要时间建立标签。', '先看点击率、完播率、加购率。', '同款做 3–5 个素材版本轮测。'], BLUE),
        ('Q2：要不要先投广告？', ['可以，但只给有自然点击苗头的素材加热。', '别拿广告救一个内容没吸引力的产品。', '先设止损线。'], GREEN),
        ('Q3：多久扩品扩站？', ['至少先跑出 1 个可复制 SKU。', '单店 SOP 稳定后再扩国家。', '没后台数据就扩张，通常死得更快。'], PINK)
    ], 17
)

add_two_col_slide(
    '恶意退货 / 投诉应对', '你不是要“情绪化应对”，而是要有证据链。',
    '恶意退货 SOP', [
        '高风险商品打包全程留档。',
        '收回退货时开箱录像、称重、拍细节。',
        '把“可二次销售 / 不可二次销售”标准写清。',
        '证据齐全后再走平台申诉。'
    ],
    '恶意投诉 SOP', [
        '先分类型：侵权、虚假宣传、竞品举报、物流时效。',
        '侵权类先下架核查，不要嘴硬。',
        '原创素材、授权书、供货链证明越齐全，胜率越高。',
        '沉淀内部黑名单与证据素材库。'
    ],
    18, left_accent=PINK, right_accent=BLUE
)

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_header(slide, '建议时间表', '别把两个月的事幻想成两周完成。', 19)
rows = [
    ['阶段', '目标', '关键动作', '通过标准'],
    ['第1–2周', '确定路径', '身份核查、主体准备、市场选择、预算表', '知道先做哪条路'],
    ['第3–6周', '完成准备', '开主体/收款/物流方案/素材规划', '具备提交入驻条件'],
    ['第2–3月', '冷启动', '上 SKU、拍素材、测广告、跑客服 SOP', '出现首批稳定订单'],
    ['第3–4月', '优化模型', '砍差品、补高转化品、优化履约', '毛利结构更清晰'],
    ['第4月以后', '放大复制', '扩站点、扩达人、扩库存', '在利润基础上扩张']
]
add_matrix_table(slide, Inches(0.62), Inches(1.28), [Inches(1.7), Inches(2.0), Inches(4.5), Inches(3.95)], rows, row_h=0.8)
add_note_box(slide, Inches(0.62), Inches(6.02), Inches(12.1), Inches(0.52), '提醒', ['多数新手不是死于“没机会”，而是死于太着急。'], YELLOW_BG, YELLOW_TXT)
add_footer(slide)

add_two_col_slide(
    '双向机遇：SEA / 反向海淘', '韩国主体的价值，不只是一条单向卖货链路。',
    '方向1：韩国 → SEA', [
        '卖的是韩系审美 + 韩国主体背书 + 内容表达。',
        '适合家居、配件、宠物、小礼品、非功效型美妆工具。',
        '重点是多语言内容和跨境履约，而不是压大库存。'
    ],
    '方向2：韩国 → 中国反向海淘', [
        '适合小众品牌、线下发现、买手式内容。',
        '玩法是内容种草 + 预售 / 代发，不一定依赖 TikTok Shop。',
        '中韩语言与文化差更容易做出差异化内容。'
    ],
    20, left_accent=GREEN, right_accent=BLUE
)

add_two_col_slide(
    'Amazon为何是补充路径', 'TikTok Shop 和 Amazon 底层逻辑不一样，新手别同时重度启动两套系统。',
    'TikTok Shop 逻辑', [
        '内容种草 + 短视频转化 + 快速测品。',
        '更适合通过素材和达人测试商品反应。',
        '节奏更快，但平台波动和运营执行也更敏感。'
    ],
    'Amazon 逻辑', [
        '货架搜索 + Listing 优化 + 评论积累 + 广告运营。',
        '更适合有稳定供应链、有备货与广告承受力的人。',
        '如果你擅长搜索电商和长期品牌化，Amazon 更有意义。'
    ],
    21, note=('核心观点', ['新手同时重度做两边，最常见结果是两边都做不好。'], (PINK_BG, PINK_TXT))
)

add_three_col_slide(
    'Amazon站点对比：美国 / 英国 / 日本', '日本站通常更适合韩国主体优先评估。', [
        ('美国站', ['市场最大，竞争最激烈，广告 CPC 最高。', '门槛高，FBA 费用贵，退货压力大。', '适合大体量、长期品牌化。'], PINK),
        ('英国站', ['英语市场，规则成熟。', 'VAT / EPR 合规要求突出。', '整体门槛中高。'], BLUE),
        ('日本站（推荐）', ['离韩国近，物流与时差友好。', '审美与韩国更接近。', '广告压力低于美国，但需日语本地化。'], GREEN)
    ], 22, note=('判断', ['韩国主体如果要把 Amazon 当补充路径，优先评估日本站。'], (GREEN_BG, GREEN_TXT))
)

add_two_col_slide(
    'Amazon注册流程与资质', '亚马逊最怕的不是材料多，而是信息不一致。',
    '注册流程', [
        '确定站点 → 准备公司资料 → 法人身份信息。',
        '准备企业地址账单 → KYC 审核。',
        '完成收款与税务设置 → 上架 SKU。'
    ],
    '所需资料与踩坑', [
        '营业执照、法人证件、企业地址证明、企业银行账户、受益人信息。',
        '主体不稳定、信息不一致、资料临时拼凑，是最常见的审核雷区。',
        '别一边注册一边改资料，容易把自己送进反复审核。'
    ],
    23
)

add_two_col_slide(
    'Amazon费用结构', '亚马逊不是“上架就卖”，而是处处都在收你成本。',
    '主要费用', [
        '月租费：US $39.99 / UK £25 / JP ¥4900。',
        '销售佣金：通常 8%–15%，按类目变化。',
        'FBA：仓储 + 拣货 + 包装 + 配送。',
        '广告 CPC：美国通常最高。'
    ],
    '隐藏成本', [
        '退货损耗、头程、包材、汇率、税务代理。',
        '利润公式：到手利润 = 售价 - 佣金 - FBA - 广告 - 头程 - 包材 - 退货损耗 - 税务。',
        '只看销量不看净利润，是亚马逊新手经典死法。'
    ],
    24, right_accent=PINK
)

add_two_col_slide(
    'Amazon税务与合规', '平台帮你代扣一部分，不等于你的合规工作结束。',
    '主要差异', [
        '美国：Sales Tax，部分州平台代收代缴。',
        '英国：VAT，需要注册与申报。',
        '欧洲：EPR，包装/电子/电池等都可能涉及。',
        '日本：消费税 + 进口报关。'
    ],
    '底线要求', [
        '平台代扣 ≠ 合规结束。',
        '必须保存完整台账、发票、物流单据与报税底稿。',
        '一旦主体、仓储、进口申报方式变化，合规口径也可能跟着变。'
    ],
    25, note=('记住', ['税务从来不是“后面再补”，而是开卖前就要设计。'], (YELLOW_BG, YELLOW_TXT))
)

add_two_col_slide(
    'Amazon热门品类 / 避坑 / 新手建议', '别把 Amazon 当无脑铺货平台。',
    '适合 / 不适合', [
        '适合新手：家居收纳、厨房小工具、桌面整理、宠物用品、低监管配件。',
        '不适合：化妆品、食品、儿童用品、带电产品、易碎高退货、侵权商品。',
        '类目越复杂，售后、认证和广告坑越深。'
    ],
    '起步建议', [
        '不要铺太多 SKU，不要只看销量不看利润，不要碰高风险类目。',
        '日本站优先，3–10 个 SKU 起步。',
        '先跑通上架 → 收款 → FBA → 广告闭环，再考虑扩品。'
    ],
    26, left_accent=GREEN, right_accent=BLUE
)

add_two_col_slide(
    '学习资源 + 行动清单', '最后一页只做一件事：把下一步讲清楚。',
    '建议盯住的资源', [
        'TikTok Shop Seller Center / Seller Academy。',
        'Amazon Sell US / UK / JP 官方定价与政策页。',
        '韩国税务师 / 会计事务所。',
        'MFDS、韩国国税厅、目标市场海关与税务官网。'
    ],
    '现在就执行的 5 步', [
        '先选路径：韩国本地 / SEA 跨境 / 反向海淘 / Amazon 补充。',
        '砍掉高监管品类，留下 3 个可测方向。',
        '用预算表重算佣金、物流、汇损、退货。',
        '先做 MVP 店铺，而不是豪华版店铺。',
        '优先跑通一条闭环，再考虑第二平台。'
    ],
    27, note=('最终结论', ['当前最现实的 0→1 路线，仍然是“韩国主体 + SEA 跨境 + 低监管品类 + 小步快跑”。'], (GREEN_BG, GREEN_TXT))
)

prs.save(OUTPUT)
print(OUTPUT)
